"""Text extraction for every supported file type.

Heavy libraries (PyMuPDF, python-docx, openpyxl, python-pptx, pytesseract) are
imported lazily so the app starts even if an optional dependency is missing —
those files are simply skipped with a recorded reason.
"""
from __future__ import annotations

from pathlib import Path

from app.utils.logging import get_logger

log = get_logger(__name__)

CODE_EXTS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".json", ".yaml", ".yml", ".toml",
    ".sh", ".zsh", ".sql", ".html", ".css", ".c", ".cpp", ".h", ".java",
    ".go", ".rs", ".rb", ".swift", ".kt",
}
TEXT_EXTS = {".txt", ".md", ".rst", ".csv", ".log", ".ini", ".cfg", ".env.example"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".heic"}
DOC_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}

SUPPORTED_EXTS = CODE_EXTS | TEXT_EXTS | IMAGE_EXTS | DOC_EXTS


def classify(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in CODE_EXTS:
        return "code"
    if ext in {".pdf", ".docx", ".pptx", ".txt", ".md", ".rst"}:
        return "document"
    if ext in {".xlsx", ".csv", ".json"}:
        return "data"
    return "other"


class ExtractionError(Exception):
    pass


def extract_text(path: Path, ocr_enabled: bool = True, max_chars: int = 200_000) -> str:
    """Return plain text for a file, or raise ExtractionError."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            text = _extract_pdf(path)
        elif ext == ".docx":
            text = _extract_docx(path)
        elif ext == ".xlsx":
            text = _extract_xlsx(path)
        elif ext == ".pptx":
            text = _extract_pptx(path)
        elif ext in IMAGE_EXTS:
            if not ocr_enabled:
                raise ExtractionError("OCR disabled")
            text = _extract_image_ocr(path)
        elif ext in CODE_EXTS or ext in TEXT_EXTS:
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            raise ExtractionError(f"unsupported extension {ext}")
    except ExtractionError:
        raise
    except Exception as e:  # noqa: BLE001 — any parser failure becomes a skip reason
        raise ExtractionError(f"{type(e).__name__}: {e}") from e
    return text[:max_chars].strip()


def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF

    parts: list[str] = []
    with fitz.open(path) as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in wb.worksheets:
            parts.append(f"# Sheet: {sheet.title}")
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i > 500:
                    break
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _find_tesseract() -> str | None:
    """Locate tesseract even when PATH lacks Homebrew dirs (launchd)."""
    import shutil

    found = shutil.which("tesseract")
    if found:
        return found
    for candidate in ("/opt/homebrew/bin/tesseract", "/usr/local/bin/tesseract"):
        if Path(candidate).exists():
            return candidate
    return None


_ocr_langs_cache: str | None = None


def _ocr_langs(binary: str) -> str:
    """Use eng+hin when the Hindi pack is installed, else eng. Cached."""
    global _ocr_langs_cache
    if _ocr_langs_cache is not None:
        return _ocr_langs_cache
    langs = "eng"
    try:
        import subprocess

        out = subprocess.run(
            [binary, "--list-langs"], capture_output=True, text=True, timeout=10
        ).stdout
        if "hin" in out.split():
            langs = "eng+hin"
    except (OSError, subprocess.SubprocessError):
        pass
    _ocr_langs_cache = langs
    return langs


def _extract_image_ocr(path: Path) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as e:
        raise ExtractionError(f"OCR deps missing: {e}") from e
    binary = _find_tesseract()
    if not binary:
        raise ExtractionError("tesseract binary not installed (brew install tesseract)")
    pytesseract.pytesseract.tesseract_cmd = binary
    try:
        with Image.open(path) as img:
            text = pytesseract.image_to_string(img, lang=_ocr_langs(binary))
    except pytesseract.TesseractNotFoundError as e:
        raise ExtractionError("tesseract binary not installed (brew install tesseract)") from e
    if not text.strip():
        raise ExtractionError("no text found in image")
    return text


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 150) -> list[str]:
    """Split text into overlapping chunks on paragraph boundaries where possible."""
    if len(text) <= chunk_size:
        return [text] if text.strip() else []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            # prefer to break at a newline near the end of the window
            nl = text.rfind("\n", start + chunk_size // 2, end)
            if nl != -1:
                end = nl
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks
